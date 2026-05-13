#!/usr/bin/env python3
"""Generate a clean, minimal flowchart slide."""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
import os

# ── Colours ──
HEAT_GREEN = RGBColor(0x00, 0x7A, 0x33)
WHITE      = RGBColor(0xFF, 0xFF, 0xFF)
BLACK      = RGBColor(0x00, 0x00, 0x00)
DARK_GRAY  = RGBColor(0x33, 0x33, 0x33)
MID_GRAY   = RGBColor(0x55, 0x55, 0x55)
LIGHT_GRAY = RGBColor(0x99, 0x99, 0x99)
ARROW_CLR  = RGBColor(0xAA, 0xAA, 0xAA)
LOOP_CLR   = RGBColor(0xD9, 0x53, 0x19)

# Step colours (muted, professional)
COLORS = [
    RGBColor(0x5B, 0x21, 0xB6),  # 1 purple
    RGBColor(0x1E, 0x40, 0xAF),  # 2 blue
    RGBColor(0x06, 0x5F, 0x46),  # 3 teal
    RGBColor(0x92, 0x40, 0x0E),  # 4 amber
    RGBColor(0x00, 0x7A, 0x33),  # 5 green
    RGBColor(0xB9, 0x18, 0x18),  # 6 red
    RGBColor(0x1E, 0x3A, 0x8A),  # 7 navy
    RGBColor(0x7C, 0x3A, 0xED),  # 8 violet
]

SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)

prs = Presentation()
prs.slide_width = SLIDE_W
prs.slide_height = SLIDE_H
slide = prs.slides.add_slide(prs.slide_layouts[6])

# White background
slide.background.fill.solid()
slide.background.fill.fore_color.rgb = WHITE

# ── Top accent bar ──
bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, SLIDE_W, Inches(0.06))
bar.fill.solid(); bar.fill.fore_color.rgb = HEAT_GREEN; bar.line.fill.background()

# ── Title ──
tx = slide.shapes.add_textbox(Inches(0.6), Inches(0.3), Inches(10), Inches(0.6))
tf = tx.text_frame; tf.word_wrap = True
p = tf.paragraphs[0]
p.text = "FROM IDEA TO LIVE DASHBOARD"; p.font.size = Pt(30)
p.font.color.rgb = HEAT_GREEN; p.font.bold = True; p.font.name = 'Arial'

# ── Subtitle ──
tx = slide.shapes.add_textbox(Inches(0.6), Inches(0.9), Inches(10), Inches(0.4))
tf = tx.text_frame; tf.word_wrap = True
p = tf.paragraphs[0]
p.text = "Step-by-step workflow and tools used at each stage"
p.font.size = Pt(16); p.font.color.rgb = MID_GRAY; p.font.name = 'Arial'; p.font.italic = True

# ── Steps data ──
steps = [
    ("Define Vision",        "Your domain expertise\nPartner reports & data"),
    ("Set Up Tools",         "Claude Max  ·  GitHub\nSupabase  ·  Cloudflare"),
    ("Design Architecture",  "Claude Code (CLI)\nProposes stack & structure"),
    ("Prepare Data",         "Excel / CSV  ·  Supabase\nClaude Code transforms"),
    ("Build Components",     "Claude Code  ·  SvelteKit\nECharts  ·  D3.js"),
    ("Test & Iterate",       "Browser (localhost)\nClaude Code fixes errors"),
    ("Deploy to Web",        "Git push → GitHub\n→ Cloudflare auto-deploys"),
    ("Feedback & Improve",   "Stakeholder input\n→ Back to step 5"),
]

# ── Layout ──
box_w = Inches(2.55)
box_h = Inches(2.05)
gap   = Inches(0.45)
start_x = Inches(0.65)
row1_y = Inches(1.65)
row2_y = Inches(4.35)


def draw_box(slide, x, y, idx, title, tools, color):
    """One clean step box: colored header strip + white body."""

    # ── White card ──
    card = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, y, box_w, box_h)
    card.fill.solid(); card.fill.fore_color.rgb = WHITE
    card.line.color.rgb = RGBColor(0xDD, 0xDD, 0xDD); card.line.width = Pt(0.75)
    card.shadow.inherit = False

    # ── Colored header band ──
    hdr_h = Inches(0.6)
    hdr = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, y, box_w, hdr_h)
    hdr.fill.solid(); hdr.fill.fore_color.rgb = color
    hdr.line.fill.background(); hdr.shadow.inherit = False
    # Mask bottom corners with a small rect so only top is rounded
    mask = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, y + Inches(0.25), box_w, Inches(0.35))
    mask.fill.solid(); mask.fill.fore_color.rgb = color
    mask.line.fill.background(); mask.shadow.inherit = False

    # ── Step number + title in header ──
    tx = slide.shapes.add_textbox(x + Inches(0.15), y + Inches(0.1), box_w - Inches(0.3), hdr_h - Inches(0.15))
    tf = tx.text_frame; tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = f"{idx}.  {title}"; p.font.size = Pt(13)
    p.font.color.rgb = WHITE; p.font.bold = True; p.font.name = 'Arial'
    p.alignment = PP_ALIGN.LEFT

    # ── Tools text below header ──
    tx = slide.shapes.add_textbox(x + Inches(0.2), y + hdr_h + Inches(0.15),
                                   box_w - Inches(0.4), box_h - hdr_h - Inches(0.2))
    tf = tx.text_frame; tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = tools; p.font.size = Pt(10.5)
    p.font.color.rgb = MID_GRAY; p.font.name = 'Arial'; p.line_spacing = 1.4


def draw_arrow_h(slide, x, y):
    """Horizontal arrow between boxes."""
    mid_y = y + box_h / 2
    shaft_w = gap - Inches(0.15)
    sh = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE,
        x + box_w + Inches(0.05), mid_y - Inches(0.02), shaft_w, Inches(0.04))
    sh.fill.solid(); sh.fill.fore_color.rgb = ARROW_CLR; sh.line.fill.background()
    # Arrowhead
    hd = slide.shapes.add_shape(MSO_SHAPE.ISOSCELES_TRIANGLE,
        x + box_w + shaft_w, mid_y - Inches(0.09), Inches(0.16), Inches(0.18))
    hd.fill.solid(); hd.fill.fore_color.rgb = ARROW_CLR
    hd.line.fill.background(); hd.rotation = 90.0


# ── Draw row 1 ──
for i in range(4):
    x = start_x + i * (box_w + gap)
    draw_box(slide, x, row1_y, i + 1, steps[i][0], steps[i][1], COLORS[i])
    if i < 3:
        draw_arrow_h(slide, x, row1_y)

# ── Draw row 2 ──
for i in range(4):
    x = start_x + i * (box_w + gap)
    draw_box(slide, x, row2_y, i + 5, steps[i + 4][0], steps[i + 4][1], COLORS[i + 4])
    if i < 3:
        draw_arrow_h(slide, x, row2_y)

# ── U-turn connector: row 1 end → row 2 start ──
turn_x = start_x + 4 * (box_w + gap) - Inches(0.15)
r1_mid = row1_y + box_h / 2
r2_mid = row2_y + box_h / 2

# Horizontal stub out of step 4
sh = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE,
    start_x + 3 * (box_w + gap) + box_w, r1_mid - Inches(0.02),
    turn_x - (start_x + 3 * (box_w + gap) + box_w), Inches(0.04))
sh.fill.solid(); sh.fill.fore_color.rgb = ARROW_CLR; sh.line.fill.background()

# Vertical drop
sh = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE,
    turn_x - Inches(0.02), r1_mid, Inches(0.04), r2_mid - r1_mid)
sh.fill.solid(); sh.fill.fore_color.rgb = ARROW_CLR; sh.line.fill.background()

# Horizontal stub into step 8
sh = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE,
    start_x + 3 * (box_w + gap) + box_w, r2_mid - Inches(0.02),
    turn_x - (start_x + 3 * (box_w + gap) + box_w), Inches(0.04))
sh.fill.solid(); sh.fill.fore_color.rgb = ARROW_CLR; sh.line.fill.background()

# Corner dots
for cy in [r1_mid, r2_mid]:
    d = slide.shapes.add_shape(MSO_SHAPE.OVAL,
        turn_x - Inches(0.06), cy - Inches(0.06), Inches(0.12), Inches(0.12))
    d.fill.solid(); d.fill.fore_color.rgb = ARROW_CLR; d.line.fill.background()

# ── Feedback loop: step 8 → step 5 (bottom arc) ──
loop_y = row2_y + box_h + Inches(0.2)
s5_mid = start_x + box_w / 2
s8_mid = start_x + 3 * (box_w + gap) + box_w / 2

# Horizontal bar
sh = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE,
    s5_mid, loop_y - Inches(0.02), s8_mid - s5_mid, Inches(0.04))
sh.fill.solid(); sh.fill.fore_color.rgb = LOOP_CLR; sh.line.fill.background()

# Vertical from step 8 down
sh = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE,
    s8_mid - Inches(0.02), row2_y + box_h, Inches(0.04), loop_y - row2_y - box_h)
sh.fill.solid(); sh.fill.fore_color.rgb = LOOP_CLR; sh.line.fill.background()

# Arrow up into step 5
sh = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE,
    s5_mid - Inches(0.02), loop_y, Inches(0.04), Inches(0.001))
sh.fill.solid(); sh.fill.fore_color.rgb = LOOP_CLR; sh.line.fill.background()

hd = slide.shapes.add_shape(MSO_SHAPE.ISOSCELES_TRIANGLE,
    s5_mid - Inches(0.09), row2_y + box_h, Inches(0.18), loop_y - row2_y - box_h + Inches(0.02))
hd.fill.solid(); hd.fill.fore_color.rgb = LOOP_CLR
hd.line.fill.background(); hd.rotation = 0.0  # pointing up

# Loop label
tx = slide.shapes.add_textbox(s5_mid + Inches(0.4), loop_y - Inches(0.25), Inches(4), Inches(0.22))
tf = tx.text_frame; p = tf.paragraphs[0]
p.text = "Iterate — stakeholder feedback drives improvements"
p.font.size = Pt(9); p.font.color.rgb = LOOP_CLR
p.font.bold = True; p.font.italic = True; p.font.name = 'Arial'

# ── Footer ──
tx = slide.shapes.add_textbox(Inches(0.6), Inches(7.0), Inches(3), Inches(0.3))
tf = tx.text_frame; p = tf.paragraphs[0]
p.text = "AI deployment strategy"; p.font.size = Pt(9)
p.font.color.rgb = LIGHT_GRAY; p.font.name = 'Arial'

tx = slide.shapes.add_textbox(Inches(10.8), Inches(6.95), Inches(1.6), Inches(0.4))
tf = tx.text_frame; p = tf.paragraphs[0]
p.text = "HEAT"; p.font.size = Pt(22); p.font.color.rgb = HEAT_GREEN
p.font.bold = True; p.font.name = 'Arial'; p.alignment = PP_ALIGN.RIGHT

# ── Save ──
out_dir = os.path.dirname(os.path.abspath(__file__))
out_path = os.path.join(out_dir, "CoolProgress_Flowchart_v2.pptx")
prs.save(out_path)
print(f"Saved: {out_path}")
