#!/usr/bin/env python3
"""Generate a 2-slide PowerPoint presentation for the CoolProgress Dashboard."""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.oxml.ns import qn
import math
import os

# ── Colours ──
DARK_BG      = RGBColor(0x0C, 0x1E, 0x3A)
DARK_BG2     = RGBColor(0x0F, 0x17, 0x2A)
BLUE_700     = RGBColor(0x1A, 0x3F, 0x7A)
BLUE_500     = RGBColor(0x25, 0x63, 0xEB)
BLUE_400     = RGBColor(0x60, 0xA5, 0xFA)
BLUE_300     = RGBColor(0x93, 0xC5, 0xFD)
CYAN_400     = RGBColor(0x22, 0xD3, 0xEE)
CYAN_300     = RGBColor(0x67, 0xE8, 0xF9)
GREEN_400    = RGBColor(0x4A, 0xDE, 0x80)
RED_400      = RGBColor(0xF8, 0x71, 0x71)
PURPLE_500   = RGBColor(0x7C, 0x3A, 0xED)
PURPLE_300   = RGBColor(0xC0, 0x84, 0xFC)
PURPLE_BG    = RGBColor(0x1E, 0x1B, 0x4B)
WHITE        = RGBColor(0xFF, 0xFF, 0xFF)
GRAY_100     = RGBColor(0xF3, 0xF4, 0xF6)
GRAY_200     = RGBColor(0xE5, 0xE7, 0xEB)
GRAY_300     = RGBColor(0xD1, 0xD5, 0xDB)
GRAY_400     = RGBColor(0x9C, 0xA3, 0xAF)
GRAY_500     = RGBColor(0x6B, 0x72, 0x80)
GRAY_600     = RGBColor(0x4B, 0x55, 0x63)
CARD_BG      = RGBColor(0x16, 0x2A, 0x50)
CARD_BG2     = RGBColor(0x1E, 0x24, 0x4A)

SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)

prs = Presentation()
prs.slide_width  = SLIDE_W
prs.slide_height = SLIDE_H

# ── Helpers ──

def set_slide_bg(slide, color):
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color

def add_gradient_bg(slide, c1, c2):
    """Set a two-stop gradient background on a slide."""
    bg = slide.background
    fill = bg.fill
    fill.gradient()
    fill.gradient_stops[0].color.rgb = c1
    fill.gradient_stops[0].position = 0.0
    fill.gradient_stops[1].color.rgb = c2
    fill.gradient_stops[1].position = 1.0

def add_rect(slide, left, top, width, height, fill_color, border_color=None, border_width=Pt(0), corner_radius=None):
    shape = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE if corner_radius else MSO_SHAPE.RECTANGLE,
        left, top, width, height
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    if border_color:
        shape.line.color.rgb = border_color
        shape.line.width = border_width
    else:
        shape.line.fill.background()
    shape.shadow.inherit = False
    # No text frame rotation
    return shape

def add_textbox(slide, left, top, width, height, text, font_size=12,
                color=WHITE, bold=False, alignment=PP_ALIGN.LEFT, font_name='Calibri'):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.color.rgb = color
    p.font.bold = bold
    p.font.name = font_name
    p.alignment = alignment
    return txBox

def add_rich_textbox(slide, left, top, width, height, runs_list, alignment=PP_ALIGN.LEFT,
                     line_spacing=1.15):
    """runs_list = [(text, font_size, color, bold, font_name), ...]"""
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = alignment
    p.line_spacing = line_spacing
    for i, (text, fs, clr, bld, fn) in enumerate(runs_list):
        if i == 0:
            run = p.runs[0] if p.runs else p.add_run()
            run.text = text
        else:
            run = p.add_run()
            run.text = text
        run.font.size = Pt(fs)
        run.font.color.rgb = clr
        run.font.bold = bld
        run.font.name = fn
    return txBox

def add_stat_card(slide, left, top, number, label, accent=CYAN_400):
    w, h = Inches(2.55), Inches(1.35)
    card = add_rect(slide, left, top, w, h, CARD_BG, border_color=BLUE_700, border_width=Pt(1), corner_radius=True)
    add_textbox(slide, left + Inches(0.2), top + Inches(0.15), Inches(2.1), Inches(0.55),
                number, font_size=30, color=accent, bold=True)
    add_textbox(slide, left + Inches(0.2), top + Inches(0.65), Inches(2.1), Inches(0.65),
                label, font_size=9.5, color=GRAY_300)

def add_pillar_card(slide, left, top, icon, name):
    w, h = Inches(1.2), Inches(0.85)
    card = add_rect(slide, left, top, w, h, RGBColor(0x12, 0x22, 0x44),
                    border_color=RGBColor(0x1E, 0x3A, 0x5F), border_width=Pt(0.75), corner_radius=True)
    add_textbox(slide, left, top + Inches(0.08), w, Inches(0.4),
                icon, font_size=18, alignment=PP_ALIGN.CENTER)
    add_textbox(slide, left, top + Inches(0.48), w, Inches(0.3),
                name, font_size=7.5, color=GRAY_300, bold=True, alignment=PP_ALIGN.CENTER)

def add_kpi_mini(slide, left, top, label, value, value_color):
    w, h = Inches(1.6), Inches(0.85)
    card = add_rect(slide, left, top, w, h, RGBColor(0x12, 0x22, 0x44),
                    border_color=RGBColor(0x1E, 0x3A, 0x5F), border_width=Pt(0.75), corner_radius=True)
    add_textbox(slide, left + Inches(0.1), top + Inches(0.08), Inches(1.4), Inches(0.3),
                label, font_size=7, color=GRAY_400, bold=True, alignment=PP_ALIGN.CENTER)
    add_textbox(slide, left + Inches(0.1), top + Inches(0.35), Inches(1.4), Inches(0.45),
                value, font_size=18, color=value_color, bold=True, alignment=PP_ALIGN.CENTER)


# ═══════════════════════════════════════════════════════
#  SLIDE 1: The CoolProgress Dashboard
# ═══════════════════════════════════════════════════════

slide1 = prs.slides.add_slide(prs.slide_layouts[6])  # blank
add_gradient_bg(slide1, DARK_BG, BLUE_700)

# ── Decorative accent bar ──
add_rect(slide1, Inches(0), Inches(0), Inches(0.08), SLIDE_H, CYAN_400)

# ── TAG ──
tag_shape = add_rect(slide1, Inches(0.6), Inches(0.55), Inches(1.75), Inches(0.35),
                     RGBColor(0x0E, 0x2B, 0x52), border_color=CYAN_400, border_width=Pt(0.75), corner_radius=True)
add_textbox(slide1, Inches(0.6), Inches(0.57), Inches(1.75), Inches(0.35),
            "AI-BUILT DASHBOARD", font_size=9, color=CYAN_300, bold=True, alignment=PP_ALIGN.CENTER)

# ── Title ──
add_textbox(slide1, Inches(0.6), Inches(1.1), Inches(5.8), Inches(0.75),
            "CoolProgress Dashboard", font_size=36, color=WHITE, bold=True)

# ── Subtitle ──
add_textbox(slide1, Inches(0.6), Inches(1.85), Inches(5.5), Inches(0.75),
            "A full-stack data platform tracking the global cooling transition — built entirely with AI assistance in a fraction of traditional development time.",
            font_size=13, color=BLUE_300)

# ── Four stat cards (2x2) ──
col1 = Inches(0.6)
col2 = Inches(3.3)
row1 = Inches(2.9)
row2 = Inches(4.4)

add_stat_card(slide1, col1, row1, "7", "Interactive data pillars\n(Emissions, MEPS, Kigali, Access,\nPolicy, Partners, Overview)", CYAN_400)
add_stat_card(slide1, col2, row1, "10+", "Global partner data sources\nintegrated (CLASP, IEA,\nSE4ALL, UNEP...)", CYAN_400)
add_stat_card(slide1, col1, row2, "150K+", "Data records powering\ninteractive maps, charts,\nand country profiles", GREEN_400)
add_stat_card(slide1, col2, row2, "190+", "Countries covered with\nfiltering by region, scenario,\nand year", GREEN_400)

# ── Tech badges ──
techs = ["SvelteKit", "TypeScript", "ECharts", "D3.js", "Supabase", "Cloudflare"]
x_start = Inches(0.6)
for i, tech in enumerate(techs):
    bw = Inches(1.15)
    bx = x_start + i * Inches(1.0)
    add_rect(slide1, bx, Inches(5.95), Inches(0.95), Inches(0.3),
             RGBColor(0x12, 0x22, 0x44), border_color=RGBColor(0x2A, 0x4A, 0x70),
             border_width=Pt(0.5), corner_radius=True)
    add_textbox(slide1, bx, Inches(5.95), Inches(0.95), Inches(0.3),
                tech, font_size=8, color=GRAY_300, bold=True, alignment=PP_ALIGN.CENTER)

# ── RIGHT SIDE: Dashboard mockup ──
mockup_left = Inches(6.6)
mockup_top = Inches(0.5)
mockup_w = Inches(6.3)
mockup_h = Inches(6.55)

# Outer mockup frame
add_rect(slide1, mockup_left, mockup_top, mockup_w, mockup_h,
         RGBColor(0x0A, 0x16, 0x30), border_color=BLUE_700, border_width=Pt(1.5), corner_radius=True)

# Title bar
add_rect(slide1, mockup_left, mockup_top, mockup_w, Inches(0.4),
         RGBColor(0x08, 0x10, 0x24), border_color=RGBColor(0x1A, 0x2E, 0x50), border_width=Pt(0.5))
# Dots
for di, dc in enumerate([RGBColor(0xEF,0x44,0x44), RGBColor(0xFB,0xBF,0x24), RGBColor(0x22,0xC5,0x5E)]):
    dot = slide1.shapes.add_shape(MSO_SHAPE.OVAL,
        mockup_left + Inches(0.2) + di * Inches(0.22), mockup_top + Inches(0.12),
        Inches(0.15), Inches(0.15))
    dot.fill.solid()
    dot.fill.fore_color.rgb = dc
    dot.line.fill.background()

# ── 6 pillar cards inside mockup ──
pillars = [("🌍","Overview"),("🏭","Emissions"),("⚡","MEPS"),("❄️","Kigali"),("🏠","Access"),("📋","Policy")]
for i, (icon, name) in enumerate(pillars):
    px = mockup_left + Inches(0.25) + (i % 3) * Inches(1.95)
    py = mockup_top + Inches(0.6) + (i // 3) * Inches(1.0)
    add_pillar_card(slide1, px, py, icon, name)

# ── Mini line chart area ──
chart_left = mockup_left + Inches(0.25)
chart_top  = mockup_top + Inches(2.75)
chart_w    = Inches(5.8)
chart_h    = Inches(1.6)
add_rect(slide1, chart_left, chart_top, chart_w, chart_h,
         RGBColor(0x0E, 0x1A, 0x36), border_color=RGBColor(0x1A, 0x2E, 0x50),
         border_width=Pt(0.5), corner_radius=True)
add_textbox(slide1, chart_left + Inches(0.15), chart_top + Inches(0.1), Inches(3), Inches(0.25),
            "GLOBAL AC STOCK PROJECTION (BILLIONS)", font_size=7.5, color=GRAY_400, bold=True)

# Draw chart line with freeform (approximate with connected shapes)
# Instead, draw data point markers and labels
data_points = [
    (0.3,  1.25, "1990", "0.6B"),
    (1.5,  1.15, "2000", "0.8B"),
    (2.9,  0.85, "2020", "2.0B"),
    (4.3,  0.45, "2035", "4.2B"),
    (5.4,  0.2,  "2050", "6.5B"),
]
for dx, dy, yr, val in data_points:
    dot = slide1.shapes.add_shape(MSO_SHAPE.OVAL,
        chart_left + Inches(dx), chart_top + Inches(0.35 + dy),
        Inches(0.12), Inches(0.12))
    dot.fill.solid()
    dot.fill.fore_color.rgb = CYAN_400
    dot.line.fill.background()
    # Year label below
    add_textbox(slide1, chart_left + Inches(dx - 0.15), chart_top + Inches(0.5 + dy),
                Inches(0.5), Inches(0.2), yr, font_size=7, color=GRAY_400, alignment=PP_ALIGN.CENTER)

# Value label on last point
add_textbox(slide1, chart_left + Inches(5.1), chart_top + Inches(0.25),
            Inches(0.7), Inches(0.25), "6.5B", font_size=12, color=CYAN_400, bold=True)

# Draw connecting lines between points using thin rectangles (approximation)
# Better: use a freeform shape
from pptx.oxml.ns import qn as _qn
from lxml import etree

def add_line_connector(slide, x1, y1, x2, y2, color, width=Pt(2)):
    """Add a simple line between two points."""
    cx = slide.shapes.add_connector(1, x1, y1, x2, y2)  # 1 = straight
    cx.line.color.rgb = color
    cx.line.width = width
    return cx

for i in range(len(data_points) - 1):
    x1 = chart_left + Inches(data_points[i][0]) + Inches(0.06)
    y1 = chart_top + Inches(0.35 + data_points[i][1]) + Inches(0.06)
    x2 = chart_left + Inches(data_points[i+1][0]) + Inches(0.06)
    y2 = chart_top + Inches(0.35 + data_points[i+1][1]) + Inches(0.06)
    add_line_connector(slide1, x1, y1, x2, y2, CYAN_400, Pt(2))

# ── 3 KPI mini cards ──
kpi_y = mockup_top + Inches(4.55)
add_kpi_mini(slide1, mockup_left + Inches(0.25), kpi_y, "EMISSIONS BAU 2050", "6,009 Mt", RED_400)
add_kpi_mini(slide1, mockup_left + Inches(2.1), kpi_y, "EMISSIONS DECARB", "1,554 Mt", GREEN_400)
add_kpi_mini(slide1, mockup_left + Inches(3.95), kpi_y, "POTENTIAL CUT", "−73%", CYAN_400)

# ── Key message at bottom of mockup ──
add_textbox(slide1, mockup_left + Inches(0.25), mockup_top + Inches(5.65), Inches(5.8), Inches(0.7),
            "\"Cooling emissions could triple by 2050. Or we cut them 73%.\"",
            font_size=11, color=BLUE_300, bold=True, alignment=PP_ALIGN.CENTER)

# Footer
add_textbox(slide1, Inches(0.6), Inches(6.9), Inches(5), Inches(0.35),
            "CoolProgress Dashboard  |  coolprogress.org", font_size=9, color=GRAY_500)


# ═══════════════════════════════════════════════════════
#  SLIDE 2: AI as a Development Partner
# ═══════════════════════════════════════════════════════

slide2 = prs.slides.add_slide(prs.slide_layouts[6])
add_gradient_bg(slide2, DARK_BG2, PURPLE_BG)

# Accent bar
add_rect(slide2, Inches(0), Inches(0), Inches(0.08), SLIDE_H, PURPLE_500)

# ── TAG ──
tag2 = add_rect(slide2, Inches(0.6), Inches(0.55), Inches(1.6), Inches(0.35),
                RGBColor(0x1C, 0x17, 0x3E), border_color=PURPLE_500, border_width=Pt(0.75), corner_radius=True)
add_textbox(slide2, Inches(0.6), Inches(0.57), Inches(1.6), Inches(0.35),
            "HOW WE BUILT IT", font_size=9, color=PURPLE_300, bold=True, alignment=PP_ALIGN.CENTER)

# ── Title ──
add_textbox(slide2, Inches(0.6), Inches(1.1), Inches(5.8), Inches(0.75),
            "AI as a Development Partner", font_size=36, color=WHITE, bold=True)

# ── Subtitle ──
add_textbox(slide2, Inches(0.6), Inches(1.85), Inches(5.3), Inches(0.65),
            "AI didn't replace the team — it amplified our capability, enabling one person to build what traditionally requires a multi-disciplinary squad.",
            font_size=13, color=RGBColor(0xC4, 0xB5, 0xFD))

# ── Process Steps (left column) ──
steps = [
    ("🧠", "Architecture & Design", "AI proposed the tech stack, component structure,\nand data model based on our requirements", PURPLE_500),
    ("💻", "Full-Stack Code Generation", "SvelteKit components, Supabase queries, ECharts\nvisualizations, and Tailwind styling — all AI-assisted", BLUE_500),
    ("🗄️", "Data Pipeline & Integration", "150K+ records from 10+ partners ingested, transformed,\nand served via Supabase with pagination", CYAN_400),
    ("🚀", "Deployment & Iteration", "Deployed on Cloudflare with rapid bug-fixing\nand feature iteration through AI pair-programming", GREEN_400),
]

for i, (icon, title, desc, accent) in enumerate(steps):
    sy = Inches(2.8) + i * Inches(1.15)
    # Icon circle
    circ = slide2.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
        Inches(0.6), sy, Inches(0.55), Inches(0.55))
    circ.fill.solid()
    r, g, b = accent[0], accent[1], accent[2]
    circ.fill.fore_color.rgb = RGBColor(r // 3, g // 3, b // 3)
    circ.line.color.rgb = accent
    circ.line.width = Pt(0.75)
    add_textbox(slide2, Inches(0.6), sy + Inches(0.05), Inches(0.55), Inches(0.45),
                icon, font_size=18, alignment=PP_ALIGN.CENTER)
    # Title
    add_textbox(slide2, Inches(1.3), sy, Inches(4.5), Inches(0.3),
                title, font_size=13, color=WHITE, bold=True)
    # Description
    add_textbox(slide2, Inches(1.3), sy + Inches(0.32), Inches(4.8), Inches(0.55),
                desc, font_size=10, color=GRAY_400)

# ═══ RIGHT COLUMN ═══

# ── Development Comparison Box ──
comp_left = Inches(6.6)
comp_top = Inches(0.5)
comp_w = Inches(6.3)
comp_h = Inches(3.5)
add_rect(slide2, comp_left, comp_top, comp_w, comp_h,
         CARD_BG2, border_color=RGBColor(0x2E, 0x2B, 0x5E), border_width=Pt(1), corner_radius=True)

add_textbox(slide2, comp_left, comp_top + Inches(0.15), comp_w, Inches(0.3),
            "DEVELOPMENT COMPARISON", font_size=10, color=GRAY_400, bold=True, alignment=PP_ALIGN.CENTER)

# Legend
legend_y = comp_top + Inches(0.5)
# Traditional
add_rect(slide2, comp_left + Inches(1.4), legend_y, Inches(0.2), Inches(0.2),
         GRAY_500, corner_radius=True)
add_textbox(slide2, comp_left + Inches(1.65), legend_y, Inches(1.3), Inches(0.2),
            "Traditional Team", font_size=8.5, color=GRAY_300)
# AI
add_rect(slide2, comp_left + Inches(3.4), legend_y, Inches(0.2), Inches(0.2),
         PURPLE_500, corner_radius=True)
add_textbox(slide2, comp_left + Inches(3.65), legend_y, Inches(1.5), Inches(0.2),
            "AI-Assisted (1 Person)", font_size=8.5, color=GRAY_300)

# Comparison bars
comparisons = [
    ("Team Size",       "4–6",   1.0,  "1",     0.2,  False),
    ("Time to MVP",     "4–6 mo", 1.0, "~6 wk", 0.25, False),
    ("Specializations", "5+",    1.0,  "1",     0.2,  False),
    ("Iterations/day",  "2–3",   0.3,  "10+",   1.0,  True),
]

bar_max_w = Inches(2.2)
for i, (label, trad_val, trad_pct, ai_val, ai_pct, ai_wins) in enumerate(comparisons):
    by = comp_top + Inches(0.9) + i * Inches(0.6)
    # Label
    add_textbox(slide2, comp_left + Inches(0.3), by + Inches(0.02), Inches(1.4), Inches(0.25),
                label, font_size=10.5, color=GRAY_300)
    # Traditional bar
    tw = int(bar_max_w * trad_pct)
    tbar = add_rect(slide2, comp_left + Inches(1.8), by, Emu(tw), Inches(0.28),
                    GRAY_500, corner_radius=True)
    add_textbox(slide2, comp_left + Inches(1.85), by, Emu(max(tw, Inches(0.5))), Inches(0.28),
                trad_val, font_size=9, color=WHITE, bold=True, alignment=PP_ALIGN.RIGHT)
    # AI bar
    aw = int(bar_max_w * ai_pct)
    abar = add_rect(slide2, comp_left + Inches(1.8) + Emu(tw) + Inches(0.1), by, Emu(aw), Inches(0.28),
                    PURPLE_500, corner_radius=True)
    add_textbox(slide2, comp_left + Inches(1.8) + Emu(tw) + Inches(0.12), by, Emu(max(aw, Inches(0.5))), Inches(0.28),
                ai_val, font_size=9, color=WHITE, bold=True, alignment=PP_ALIGN.RIGHT)

# ── "What This Means for Us" Box ──
take_top = Inches(4.2)
take_h = Inches(3.1)
add_rect(slide2, comp_left, take_top, comp_w, take_h,
         RGBColor(0x1A, 0x15, 0x3A), border_color=PURPLE_500, border_width=Pt(1), corner_radius=True)

add_textbox(slide2, comp_left + Inches(0.3), take_top + Inches(0.2), Inches(4), Inches(0.3),
            "WHAT THIS MEANS FOR US", font_size=10, color=PURPLE_300, bold=True)

takeaways = [
    ("Rapid prototyping", "Turn data into stakeholder-ready\ndashboards in weeks, not months"),
    ("Domain experts lead", "Cooling specialists drive the product;\nAI handles the engineering complexity"),
    ("Scalable approach", "Same method applies to NDC tracking,\npolicy analysis, and partner tools"),
    ("Lower barrier", "No large dev teams needed to ship\nproduction-quality analytical tools"),
]

for i, (title, desc) in enumerate(takeaways):
    ty = take_top + Inches(0.6) + i * Inches(0.6)
    # Checkmark
    add_textbox(slide2, comp_left + Inches(0.3), ty, Inches(0.25), Inches(0.25),
                "✓", font_size=13, color=GREEN_400, bold=True)
    # Title + desc
    add_textbox(slide2, comp_left + Inches(0.6), ty - Inches(0.02), Inches(5.3), Inches(0.22),
                title, font_size=11.5, color=WHITE, bold=True)
    add_textbox(slide2, comp_left + Inches(0.6), ty + Inches(0.22), Inches(5.3), Inches(0.4),
                desc, font_size=9.5, color=GRAY_300)

# Footer
add_textbox(slide2, Inches(0.6), Inches(6.9), Inches(5), Inches(0.35),
            "CoolProgress Dashboard  |  Built with Claude AI + SvelteKit", font_size=9, color=GRAY_500)

# ── Save ──
out_dir = os.path.dirname(os.path.abspath(__file__))
out_path = os.path.join(out_dir, "CoolProgress_AI_Dashboard_Presentation.pptx")
prs.save(out_path)
print(f"Saved: {out_path}")
