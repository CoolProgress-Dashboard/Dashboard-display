#!/usr/bin/env python3
"""Generate a 2-slide PPTX following the HEAT GmbH template format."""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
import os

# ── Template colours (matched from screenshot) ──
HEAT_GREEN   = RGBColor(0x00, 0x7A, 0x33)  # HEAT brand green (titles, headers)
HEAT_GREEN_D = RGBColor(0x00, 0x6B, 0x2B)  # Darker green for emphasis
BLACK        = RGBColor(0x00, 0x00, 0x00)
DARK_GRAY    = RGBColor(0x33, 0x33, 0x33)
MID_GRAY     = RGBColor(0x66, 0x66, 0x66)
LIGHT_GRAY   = RGBColor(0x99, 0x99, 0x99)
WHITE        = RGBColor(0xFF, 0xFF, 0xFF)

# Card background colours (from screenshot)
CARD_GRAY    = RGBColor(0x8E, 0x94, 0xA7)  # Blue-gray
CARD_YELLOW  = RGBColor(0xE8, 0xE8, 0x00)  # Bright yellow
CARD_GREEN   = RGBColor(0x66, 0xD9, 0x00)  # Bright green

SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)

prs = Presentation()
prs.slide_width  = SLIDE_W
prs.slide_height = SLIDE_H


# ── Helpers ──

def add_textbox(slide, left, top, width, height, text, font_size=12,
                color=BLACK, bold=False, alignment=PP_ALIGN.LEFT,
                font_name='Arial', italic=False, line_spacing=1.15):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.color.rgb = color
    p.font.bold = bold
    p.font.name = font_name
    p.font.italic = italic
    p.alignment = alignment
    p.line_spacing = line_spacing
    return txBox


def add_multiline_textbox(slide, left, top, width, height, lines,
                          font_size=12, color=BLACK, bold=False,
                          alignment=PP_ALIGN.LEFT, font_name='Arial',
                          line_spacing=1.3, bullet=False):
    """Add a text box with multiple paragraphs (one per item in lines list)."""
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True

    for i, line in enumerate(lines):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()

        if bullet:
            p.text = "• " + line
        else:
            p.text = line

        p.font.size = Pt(font_size)
        p.font.color.rgb = color
        p.font.bold = bold
        p.font.name = font_name
        p.alignment = alignment
        p.line_spacing = line_spacing
        p.space_after = Pt(4)

    return txBox


def add_rect(slide, left, top, width, height, fill_color, corner_radius=None):
    shape_type = MSO_SHAPE.ROUNDED_RECTANGLE if corner_radius else MSO_SHAPE.RECTANGLE
    shape = slide.shapes.add_shape(shape_type, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    shape.line.fill.background()
    shape.shadow.inherit = False
    return shape


def add_slide_template(slide, title_text, subtitle_text, slide_number,
                       col1_header, col1_lines,
                       col2_header, col2_lines,
                       col3_header, col3_lines):
    """Create a slide following the HEAT template format."""

    # ── White background ──
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = WHITE

    # ── Green top accent line ──
    add_rect(slide, Inches(0), Inches(0), SLIDE_W, Inches(0.06), HEAT_GREEN)

    # ── Title (bold, uppercase, green) ──
    add_textbox(slide, Inches(0.6), Inches(0.25), Inches(10), Inches(0.7),
                title_text.upper(), font_size=32, color=HEAT_GREEN, bold=True,
                font_name='Arial')

    # ── Subtitle (regular, dark gray, italic) ──
    add_textbox(slide, Inches(0.6), Inches(0.95), Inches(10), Inches(0.5),
                subtitle_text, font_size=18, color=DARK_GRAY,
                font_name='Arial', italic=True)

    # ── Thin separator line ──
    add_rect(slide, Inches(0.6), Inches(1.55), Inches(12.1), Inches(0.015), RGBColor(0xDD, 0xDD, 0xDD))

    # ── Three cards ──
    card_top    = Inches(1.85)
    card_height = Inches(4.6)
    card_width  = Inches(3.75)
    card_gap    = Inches(0.35)
    card_x1     = Inches(0.6)
    card_x2     = card_x1 + card_width + card_gap
    card_x3     = card_x2 + card_width + card_gap

    cards = [
        (card_x1, CARD_GRAY,   col1_header, col1_lines),
        (card_x2, CARD_YELLOW, col2_header, col2_lines),
        (card_x3, CARD_GREEN,  col3_header, col3_lines),
    ]

    for cx, bg_color, header, lines in cards:
        # Card background
        add_rect(slide, cx, card_top, card_width, card_height, bg_color, corner_radius=True)

        # Card header (green, bold, uppercase)
        header_color = HEAT_GREEN_D if bg_color == CARD_GREEN else HEAT_GREEN
        add_textbox(slide, cx + Inches(0.25), card_top + Inches(0.2),
                    card_width - Inches(0.5), Inches(0.45),
                    header.upper(), font_size=14, color=header_color, bold=True,
                    font_name='Arial')

        # Card content (bullet points)
        content_color = WHITE if bg_color == CARD_GRAY else DARK_GRAY
        add_multiline_textbox(
            slide, cx + Inches(0.25), card_top + Inches(0.7),
            card_width - Inches(0.5), card_height - Inches(0.9),
            lines, font_size=11.5, color=content_color,
            font_name='Arial', line_spacing=1.35, bullet=True
        )

    # ── Footer ──
    footer_y = Inches(6.85)

    # Left: "AI deployment strategy"
    add_textbox(slide, Inches(0.6), footer_y, Inches(3), Inches(0.35),
                "AI deployment strategy", font_size=9, color=LIGHT_GRAY,
                font_name='Arial')

    # Right: HEAT logo text + page number
    add_textbox(slide, Inches(10.8), footer_y, Inches(1.6), Inches(0.4),
                "HEAT", font_size=22, color=HEAT_GREEN, bold=True,
                font_name='Arial', alignment=PP_ALIGN.RIGHT)

    add_textbox(slide, Inches(12.45), footer_y + Inches(0.05), Inches(0.5), Inches(0.35),
                str(slide_number), font_size=11, color=MID_GRAY,
                font_name='Arial', alignment=PP_ALIGN.RIGHT)


# ═══════════════════════════════════════════════════════
#  SLIDE 1: CoolProgress Dashboard
# ═══════════════════════════════════════════════════════

slide1 = prs.slides.add_slide(prs.slide_layouts[6])  # blank

add_slide_template(
    slide1,
    title_text="Building the Cooling Dashboard",
    subtitle_text="My workflow: from prompt to production — the tools, the process, and what I had to learn",
    slide_number=1,

    col1_header="The Tool Chain",
    col1_lines=[
        "Claude Max subscription → gives access to Claude Code, an AI coding assistant that runs in the terminal",
        "Claude Code reads, writes, and edits project files directly — I describe what I want, it writes the code",
        "GitHub → version control. Every change gets committed, so I can always roll back if something breaks",
        "Supabase → cloud database + user login. I upload partner data (CSVs), Claude writes the queries to serve it",
        "Cloudflare Pages → hosting. I push code to GitHub, Cloudflare auto-deploys the live website in minutes",
        "The chain: I prompt Claude → it writes code → I push to GitHub → website updates automatically",
    ],

    col2_header="How the Work Actually Looked",
    col2_lines=[
        "I describe features in plain language: \"add a map showing emissions by country with a year slider\"",
        "Claude generates full components — frontend, database queries, styling — I review the output",
        "When something breaks (and it does), I paste the error back to Claude and we fix it together",
        "Data prep was key: I gathered CSVs from CLASP, SE4ALL, IEA, UNEP — cleaned and uploaded to Supabase",
        "Lots of iteration: first version is never perfect. I'd say \"make the legend clearer\" or \"the map colours are wrong\" and refine",
        "Not magic — it's like working with a very fast junior developer who needs clear direction from you",
    ],

    col3_header="What I Needed to Know",
    col3_lines=[
        "Deep knowledge of the cooling sector — AI writes code, but I decide WHAT to show and WHY",
        "Ability to describe features clearly — vague prompts = vague results. Specific asks get specific code",
        "Basic sense of how web apps work (pages, data, charts) — not coding, but understanding the pieces",
        "Patience to test and iterate — first output is ~70% right, getting to 100% takes back-and-forth",
        "Know where the data lives — I sourced partner datasets, understood their structure, and told Claude how to use them",
        "Willingness to learn on the fly — I picked up Git, Supabase basics, and deployment by doing, with AI guiding me",
    ],
)


# ═══════════════════════════════════════════════════════
#  SLIDE 2: AI Development — Impact, Improvement, Potential
# ═══════════════════════════════════════════════════════

slide2 = prs.slides.add_slide(prs.slide_layouts[6])  # blank

add_slide_template(
    slide2,
    title_text="What This Means for HEAT",
    subtitle_text="Honest takeaways — what worked, what didn't, and where this could go",
    slide_number=2,

    col1_header="Impact",
    col1_lines=[
        "A live dashboard with 7 data pillars, 190+ countries, and 150K+ records — built by one person",
        "Partners like CLASP, SE4ALL, and UNEP see their data visualized in one place for the first time",
        "Weeks, not months: from idea to working prototype to deployed product in roughly 6 weeks",
        "I'm not a developer — I'm a cooling domain expert. AI let me build something I couldn't have built alone",
        "Rapid feedback loops: stakeholders ask for a change in the morning, it's live by afternoon",
    ],

    col2_header="Need for Improvement",
    col2_lines=[
        "AI doesn't understand our domain — it writes code fast but I had to catch wrong assumptions about the data constantly",
        "Debugging can be frustrating: sometimes Claude fixes one thing and breaks another. You need patience",
        "I had to learn Git, Supabase, and deployment basics along the way — steep learning curve at the start",
        "Code quality: AI-generated code works but isn't always clean. A real developer would structure it differently",
        "Documentation is thin — if someone else needs to maintain this, they'll need onboarding",
        "Data validation is still manual — I check every chart against source data myself",
    ],

    col3_header="Potential for HEAT",
    col3_lines=[
        "Any Heatie with domain knowledge could build analytical tools this way — you don't need to be a coder",
        "NDC assessments, country reports, policy trackers — same approach, different data",
        "Client deliverables: instead of static PDFs, we could ship interactive dashboards",
        "Automate the boring parts: data ingestion, indicator calculations, report formatting",
        "Faster proposals: \"we can prototype this in 2 weeks\" instead of scoping a 3-month dev project",
        "The key insight: AI is a multiplier for domain expertise, not a replacement for it",
    ],
)


# ── Save ──
out_dir = os.path.dirname(os.path.abspath(__file__))
out_path = os.path.join(out_dir, "CoolProgress_AI_Dashboard_HEAT.pptx")
prs.save(out_path)
print(f"Saved: {out_path}")
